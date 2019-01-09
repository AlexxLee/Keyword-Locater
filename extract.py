from pptx import Presentation
import os

# Get absolute path
file_path = input("Path: ")
print(file_path)
key = input("Keyword: ")
print(key)
dirs = os.listdir(file_path)

# C:\Users\94231\Documents\Python Scripts\Locator\temp\
# Print pdf file names

text_runs = []
for file in dirs:
    if file.endswith('.pptx'):
        #print(file)
        the_file = Presentation(file_path + file)
        i = 0
        for slide in the_file.slides:
            i = i+1
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        #text_runs.append(run.text)
                        if run.text.find(key) != -1:
                            text_runs.append(file+"-slide "+str(i))


for x in text_runs:
    print(x)

