from tkinter import *
from pptx import Presentation
import sys, os

# Global variables
insuf_input_errormsg = "Insufficient input. Try again."
invalid_dir_errormsg = "Missing '/' or '\\' at the end of the directory path. Try again."
results = []

# Create a dialog with an input box and a button for user inputs
root = Tk()
root.title("Keyword Locator")
# root.geometry("620x70")


def error_handling(directory, keyword):
    if (not directory) or (not keyword):
        insuf_input_label.grid(row=2, column=0)
        return True
    else:
        insuf_input_label.grid_remove()

    if (not directory.endswith("/")) and (not directory.endswith("\\")):
        invalid_dir_label.grid(row=2, column=0)
        return True
    else:
        invalid_dir_label.grid_remove()
    return False


def search():
    #print(len(results))
    directory = path_entry.get("1.0",'end-1c')
    keyword = kw_entry.get("1.0",'end-1c')

    # Error handling
    error_exist = error_handling(directory, keyword)
    if error_exist:
        print("Input error exists!")
    else:
        getList(directory, keyword)
        #print(results[0]+"qweqwe")
        # for x in results:
        #      print(x)

    win = Tk()
    win.title('Result')
    win.geometry('200x400') # Size 200, 400

        scrollbar = Scrollbar(win)
        scrollbar.pack(side=RIGHT, fill=Y)

    listbox = Listbox(win, yscrollcommand=scrollbar.set)
    for x in results:
        listbox.insert(END, x)
    results.clear()
    listbox.pack(side=LEFT, fill=BOTH)

        scrollbar.config(command=listbox.yview)

        mainloop()

def getList(directory, keyword):
    files = os.listdir(directory)
    for file in files:
        if os.path.isdir(os.path.join(os.path.abspath(directory), file)):
            subdir = os.path.join(os.path.abspath(directory), file)+"\\"
            # print(subdir)
            getList(subdir, keyword)
        if file.endswith('.pptx'):
            # print(file)
            the_file = Presentation(directory + file)
            i = 0
            for slide in the_file.slides:
                i = i+1
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.find(keyword) != -1:
                                #print(file+"-slide "+str(i))
                                results.append(file+"-slide "+str(i))




def cancel():
    print("Exit!")
    sys.exit()


# Path entry
path_label = Label(root, text="Enter the absolute path for a directory:")
path_entry = Text(root, height=1, width=50)
path_label.grid(row=0, column=0)
path_entry.grid(row=0, column=1)

# Keyword entry
kw_label = Label(root, text="Enter the keyword:")
kw_entry = Text(root, height=1, width=50)
kw_label.grid(row=1, column=0)
kw_entry.grid(row=1, column=1)

# Error message
insuf_input_label = Label(root, text=insuf_input_errormsg, fg="red")
invalid_dir_label = Label(root, text=invalid_dir_errormsg, fg="red")


# Buttons
submit_btn = Button(root, text="Search", command=search)
exit_btn = Button(root, text="Cancel", command=exit)
submit_btn.grid(row=3, column=0)
exit_btn.grid(row=3, column=1)

root.mainloop()
